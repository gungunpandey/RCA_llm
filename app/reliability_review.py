"""
AI-Generated Reliability Review — PowerPoint builder.

Produces an 8-slide management deck (+ cover) from:
  * the verified analysis bundle (insights_engine)  — numbers
  * the historical analytics rollup                 — trends / tables
  * an LLM-written narrative (with deterministic fallback)

Design principle carries over from the insights engine: every number printed on
a slide comes from the verified bundle/analytics. The LLM only writes prose
(executive summary, recommendations, decisions) and never invents figures.
"""
import io
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION


# ── palette (matches the ProdAI teal theme) ──────────────────────────────────
TEAL = RGBColor(0x33, 0xB1, 0xB0)
TEAL_DARK = RGBColor(0x2B, 0x8C, 0x8B)
PURPLE = RGBColor(0x7C, 0x6B, 0xFF)
DARK = RGBColor(0x1F, 0x2D, 0x3D)
GREY = RGBColor(0x64, 0x74, 0x8B)
RED = RGBColor(0xDC, 0x26, 0x26)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
AMBER = RGBColor(0xB4, 0x7A, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF1, 0xF6, 0xF6)

# Chart series palette (cycled for pie/doughnut/colored bars)
_PALETTE = [
    RGBColor(0x33, 0xB1, 0xB0), RGBColor(0x7C, 0x6B, 0xFF), RGBColor(0xF5, 0x9E, 0x0B),
    RGBColor(0xEF, 0x44, 0x44), RGBColor(0x10, 0xB9, 0x81), RGBColor(0xF9, 0x73, 0x16),
    RGBColor(0x8B, 0x5C, 0xF6), RGBColor(0x06, 0xB6, 0xD4), RGBColor(0xFB, 0x92, 0x3C),
    RGBColor(0x4A, 0xDE, 0x80),
]


# ── low-level helpers ────────────────────────────────────────────────────────
def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, left, top, width, height, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _text(slide, left, top, width, height, text, size=14, color=DARK,
          bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return box


def _header(slide, prs, idx, title):
    """Teal title band with a slide number badge."""
    _rect(slide, 0, 0, prs.slide_width, Inches(0.95), TEAL)
    _text(slide, Inches(0.55), Inches(0.18), prs.slide_width - Inches(1.6),
          Inches(0.6), title, size=26, color=WHITE, bold=True)
    _text(slide, prs.slide_width - Inches(1.0), Inches(0.25), Inches(0.7),
          Inches(0.5), str(idx), size=18, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)


def _bullets(slide, items, left, top, width, height, size=15, color=DARK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = f"•  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return box


def _table(slide, headers, rows, left, top, width, row_h=Inches(0.34)):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    height = row_h * n_rows
    tbl = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = str(h)
        cell.fill.solid()
        cell.fill.fore_color.rgb = TEAL
        para = cell.text_frame.paragraphs[0]
        para.runs[0].font.size = Pt(11)
        para.runs[0].font.bold = True
        para.runs[0].font.color.rgb = WHITE
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 else LIGHT
            para = cell.text_frame.paragraphs[0]
            para.runs[0].font.size = Pt(10.5)
            para.runs[0].font.color.rgb = DARK
    return tbl


def _metric_tiles(slide, metrics, left, top, total_width, prs):
    """Render exec-summary metric tiles in a row."""
    n = min(len(metrics), 4)
    if n == 0:
        return
    gap = Inches(0.2)
    tile_w = (total_width - gap * (n - 1)) / n
    tile_h = Inches(1.25)
    for i, m in enumerate(metrics[:n]):
        x = left + i * (tile_w + gap)
        accent = RED if m.get("good") is False else GREEN
        _rect(slide, x, top, tile_w, tile_h, LIGHT)
        _rect(slide, x, top, Inches(0.08), tile_h, accent)
        _text(slide, x + Inches(0.18), top + Inches(0.12), tile_w - Inches(0.3),
              Inches(0.4), m["label"].upper(), size=9.5, color=GREY, bold=True)
        _text(slide, x + Inches(0.18), top + Inches(0.45), tile_w - Inches(0.3),
              Inches(0.55), str(m["value"]), size=24, color=DARK, bold=True)
        if m.get("trend"):
            arrow = "▲" if m.get("direction") == "up" else "▼" if m.get("direction") == "down" else ""
            _text(slide, x + Inches(0.18), top + Inches(0.95), tile_w - Inches(0.3),
                  Inches(0.3), f"{arrow} {m['trend']}", size=11, color=accent, bold=True)


# ── chart + notes helpers ────────────────────────────────────────────────────
def _notes(slide, lines):
    """Push detail text into the speaker-notes pane to keep slides visual."""
    if not lines:
        return
    tf = slide.notes_slide.notes_text_frame
    tf.text = "\n".join(f"• {l}" for l in lines)


def _style_axes(chart):
    try:
        chart.category_axis.tick_labels.font.size = Pt(9)
        chart.category_axis.tick_labels.font.color.rgb = DARK
    except Exception:
        pass
    try:
        chart.value_axis.tick_labels.font.size = Pt(9)
        chart.value_axis.tick_labels.font.color.rgb = GREY
    except Exception:
        pass


def _legend(chart, position=XL_LEGEND_POSITION.BOTTOM):
    chart.has_legend = True
    chart.legend.position = position
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(10)


def _bar_chart(slide, left, top, width, height, categories, values,
               color=TEAL, colors=None, horizontal=True):
    cd = CategoryChartData()
    cd.categories = categories
    cd.add_series("Failures", values)
    ct = XL_CHART_TYPE.BAR_CLUSTERED if horizontal else XL_CHART_TYPE.COLUMN_CLUSTERED
    chart = slide.shapes.add_chart(ct, left, top, width, height, cd).chart
    chart.has_title = False
    chart.has_legend = False
    series = chart.series[0]
    if colors:
        for i, pt in enumerate(series.points):
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = colors[i % len(colors)]
    else:
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = color
    try:
        plot = chart.plots[0]
        plot.has_data_labels = True
        plot.data_labels.font.size = Pt(9)
        plot.data_labels.font.color.rgb = DARK
    except Exception:
        pass
    _style_axes(chart)
    return chart


def _line_chart(slide, left, top, width, height, categories, series_list, legend=True):
    cd = CategoryChartData()
    cd.categories = categories
    for name, vals, _c in series_list:
        cd.add_series(name, vals)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, left, top, width, height, cd).chart
    chart.has_title = False
    if legend:
        _legend(chart)
    else:
        chart.has_legend = False
    for i, (_n, _v, color) in enumerate(series_list):
        ser = chart.series[i]
        ser.format.line.color.rgb = color
        ser.format.line.width = Pt(2.25)
        ser.smooth = False
    _style_axes(chart)
    return chart


def _doughnut(slide, left, top, width, height, categories, values):
    cd = CategoryChartData()
    cd.categories = categories
    cd.add_series("Count", values)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, left, top, width, height, cd).chart
    chart.has_title = False
    _legend(chart, XL_LEGEND_POSITION.RIGHT)
    for i, pt in enumerate(chart.series[0].points):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = _PALETTE[i % len(_PALETTE)]
    try:
        plot = chart.plots[0]
        plot.has_data_labels = True
        plot.data_labels.font.size = Pt(9)
        plot.data_labels.font.color.rgb = WHITE
    except Exception:
        pass
    return chart


def _grouped_columns(slide, left, top, width, height, categories, series_list, legend=True):
    cd = CategoryChartData()
    cd.categories = categories
    for name, vals, _c in series_list:
        cd.add_series(name, vals)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, cd).chart
    chart.has_title = False
    if legend:
        _legend(chart)
    else:
        chart.has_legend = False
    for i, (_n, _v, color) in enumerate(series_list):
        ser = chart.series[i]
        ser.format.fill.solid()
        ser.format.fill.fore_color.rgb = color
    _style_axes(chart)
    return chart


# ── deterministic narrative (fallback when LLM is unavailable) ────────────────
def deterministic_narrative(bundle: dict, analytics: dict) -> dict:
    es = bundle.get("executive_summary", {}) or {}

    exec_bullets = []
    for m in es.get("metrics", []):
        line = f"{m['label']}: {m['value']}"
        if m.get("trend"):
            line += f" ({m['trend']})"
        exec_bullets.append(line)

    rca_bullets = [f"{r['category']}: {r['count']} failures"
                   for r in (analytics.get("rootCause") or [])[:6]]

    rec = list(es.get("recommended_focus", []))
    if not rec:
        rec = ["Maintain current preventive-maintenance cadence; no critical focus areas flagged."]

    mgmt = []
    for c in bundle.get("unimplemented_capa_recurrence", [])[:3]:
        mgmt.append(f"Assign owner and deadline to close CAPA #{c['capa_id']} on "
                    f"{c['machine']} — failures are recurring while it stays open.")
    if bundle.get("capa_stats", {}).get("overdue"):
        mgmt.append(f"Review {bundle['capa_stats']['overdue']} overdue CAPA(s) for re-prioritization.")
    if not mgmt:
        mgmt = ["No critical management decisions flagged for this period."]

    return {
        "executive_summary": exec_bullets or ["No significant maintenance activity in this period."],
        "root_cause_analysis": rca_bullets or ["No failure-type data recorded for this period."],
        "recommended_actions": rec,
        "management_decisions": mgmt,
    }


# ── slide builders ───────────────────────────────────────────────────────────
def _cover(prs, meta):
    slide = _blank(prs)
    _rect(slide, 0, 0, prs.slide_width, prs.slide_height, DARK)
    _rect(slide, 0, Inches(2.55), prs.slide_width, Inches(0.06), TEAL)
    _text(slide, Inches(0.8), Inches(1.5), prs.slide_width - Inches(1.6),
          Inches(1.0), "Reliability Review", size=44, color=WHITE, bold=True)
    _text(slide, Inches(0.8), Inches(2.7), prs.slide_width - Inches(1.6),
          Inches(0.6), meta["scope"], size=20, color=TEAL, bold=True)
    _text(slide, Inches(0.8), Inches(3.4), prs.slide_width - Inches(1.6),
          Inches(0.5), f"Period: {meta['period']}", size=15, color=LIGHT)
    _text(slide, Inches(0.8), Inches(3.85), prs.slide_width - Inches(1.6),
          Inches(0.5), f"Generated by ProdAI · {meta['generated']}", size=12, color=GREY)


def _slide_exec(prs, bundle, narrative):
    slide = _blank(prs)
    _header(slide, prs, 1, "Executive Summary")
    es = bundle.get("executive_summary", {}) or {}
    _metric_tiles(slide, es.get("metrics", []), Inches(0.55), Inches(1.2),
                  prs.slide_width - Inches(1.1), prs)
    bullets = narrative.get("executive_summary", [])
    # Keep the slide light: top 4 bullets on-slide, full detail in speaker notes.
    _bullets(slide, bullets[:4], Inches(0.55), Inches(2.75),
             prs.slide_width - Inches(1.1), Inches(3.8), size=15)
    _notes(slide, bullets)


def _slide_trends(prs, analytics):
    slide = _blank(prs)
    _header(slide, prs, 2, "Breakdown Trends")
    d = analytics.get("direction") or {}
    pct = d.get("pct")
    direction = d.get("direction")
    # Big delta callout (left), trend chart (right).
    delta = "—" if pct is None else f"{'+' if direction == 'up' else ''}{pct}%"
    color = RED if direction == "up" else GREEN if direction == "down" else GREY
    _text(slide, Inches(0.55), Inches(1.25), Inches(3.2), Inches(0.95),
          delta, size=46, color=color, bold=True)
    _text(slide, Inches(0.55), Inches(2.3), Inches(3.2), Inches(0.9),
          f"failures vs prior half ({d.get('recent', 0)} vs {d.get('previous', 0)})",
          size=12, color=GREY)

    trend = analytics.get("trend") or []
    if trend:
        cats = [t.get("period", "—") for t in trend]
        fails = [t.get("failures", 0) for t in trend]
        mttr = [t.get("avg_mttr", 0) for t in trend]
        _line_chart(slide, Inches(4.0), Inches(1.15), Inches(8.85), Inches(5.5),
                    cats, [("Failures", fails, TEAL), ("Avg MTTR (h)", mttr, PURPLE)])
    else:
        _text(slide, Inches(4.0), Inches(2.0), Inches(6), Inches(0.5),
              "No trend data for this period.", size=13, color=GREY)


def _slide_top_equipment(prs, analytics):
    slide = _blank(prs)
    _header(slide, prs, 3, "Top Problem Equipment")
    top = analytics.get("top")
    if top:
        _text(slide, Inches(0.55), Inches(1.15), Inches(9), Inches(0.5),
              top["equipment_name"], size=20, color=DARK, bold=True)
        _text(slide, Inches(0.55), Inches(1.7), Inches(9), Inches(0.4),
              f"{top['failure_count']} failures   ·   {top.get('criticality', '—')} criticality"
              f"   ·   Avg MTTR {top.get('avg_mttr', 0)} h", size=12, color=GREY)
    freq = (analytics.get("freq") or [])[:8]
    if freq:
        # Horizontal bars read largest-on-top, so reverse for BAR_CLUSTERED.
        cats = [f["equipment_name"] for f in freq][::-1]
        vals = [f["failure_count"] for f in freq][::-1]
        _bar_chart(slide, Inches(0.55), Inches(2.3), Inches(12.25), Inches(4.55),
                   cats, vals, color=TEAL, horizontal=True)
    else:
        _text(slide, Inches(0.55), Inches(2.4), Inches(8), Inches(0.5),
              "No failures recorded for this period.", size=13, color=GREY)


def _slide_rca(prs, analytics, narrative):
    slide = _blank(prs)
    _header(slide, prs, 4, "Root Cause Analysis")
    rc = analytics.get("rootCause") or []
    if rc:
        cats = [r["category"] for r in rc]
        vals = [r["count"] for r in rc]
        _doughnut(slide, Inches(0.55), Inches(1.25), Inches(5.7), Inches(5.4), cats, vals)
    else:
        _text(slide, Inches(0.55), Inches(1.4), Inches(5), Inches(0.5),
              "No failure-type data recorded.", size=13, color=GREY)
    bullets = narrative.get("root_cause_analysis", [])
    _bullets(slide, bullets[:5], Inches(6.6), Inches(1.4),
             prs.slide_width - Inches(7.1), Inches(5.0), size=14)
    _notes(slide, bullets)


def _slide_capa(prs, bundle):
    slide = _blank(prs)
    _header(slide, prs, 5, "CAPA Effectiveness")
    eff = bundle.get("capa_effectiveness") or []
    if eff:
        cats = [e["machine"] for e in eff]
        before = [e["before"] for e in eff]
        after = [e["after"] for e in eff]
        _grouped_columns(slide, Inches(0.55), Inches(1.3), Inches(9.0), Inches(4.4),
                         cats, [("Before CAPA", before, GREY), ("After CAPA", after, TEAL)])
        _notes(slide, [f"{e['machine']} (CAPA #{e['capa_id']}): {e['before']} -> {e['after']} "
                       f"breakdowns ({e['change_pct']}%) over {e['window_days']} days"
                       for e in eff])
    else:
        _text(slide, Inches(0.55), Inches(1.4), Inches(11), Inches(1.0),
              "No completed CAPAs with a fully-elapsed comparison window yet. "
              "Effectiveness will populate as CAPAs are closed and time passes.",
              size=14, color=GREY)
    stats = bundle.get("capa_stats") or {}
    if stats.get("total"):
        _text(slide, Inches(0.55), Inches(6.1), Inches(11), Inches(0.5),
              f"CAPA compliance this period: {stats.get('compliance_pct', 0)}%  "
              f"({stats.get('completed', 0)}/{stats.get('total', 0)} completed, "
              f"{stats.get('overdue', 0)} overdue)", size=13, color=DARK, bold=True)


def _slide_risks(prs, bundle):
    slide = _blank(prs)
    _header(slide, prs, 6, "Predicted Risks")
    _text(slide, Inches(0.55), Inches(1.1), prs.slide_width - Inches(1.1), Inches(0.4),
          "Elevated risk based on recent activity — not a guaranteed forecast.",
          size=11, color=GREY)
    risks = bundle.get("risk_alerts") or []
    if risks:
        # Bar of recent failures per at-risk asset, colored by risk level.
        cats = [r["machine"] for r in risks][::-1]
        vals = [r["failures"] for r in risks][::-1]
        colors = [(RED if r["risk_level"] == "High" else AMBER) for r in risks][::-1]
        _bar_chart(slide, Inches(0.55), Inches(1.7), Inches(8.6), Inches(4.9),
                   cats, vals, colors=colors, horizontal=True)
        # Legend chips for risk levels.
        _text(slide, Inches(9.4), Inches(1.8), Inches(3.2), Inches(0.4),
              "■ High risk", size=12, color=RED, bold=True)
        _text(slide, Inches(9.4), Inches(2.2), Inches(3.2), Inches(0.4),
              "■ Medium risk", size=12, color=AMBER, bold=True)
        _notes(slide, [f"{r['machine']}: {r['risk_level']} risk — {'; '.join(r['reasons'][:3])}"
                       for r in risks])
    else:
        _text(slide, Inches(0.55), Inches(1.8), Inches(9), Inches(0.5),
              "No elevated-risk assets detected this period.", size=14, color=GREEN, bold=True)


def _slide_actions(prs, narrative):
    slide = _blank(prs)
    _header(slide, prs, 7, "Recommended Actions")
    items = narrative.get("recommended_actions", [])
    _bullets(slide, items[:6], Inches(0.6), Inches(1.3),
             prs.slide_width - Inches(1.2), Inches(5.0), size=16)
    _notes(slide, items)


def _slide_decisions(prs, narrative):
    slide = _blank(prs)
    _header(slide, prs, 8, "Management Decisions Required")
    items = narrative.get("management_decisions", [])
    _bullets(slide, items[:6], Inches(0.6), Inches(1.3),
             prs.slide_width - Inches(1.2), Inches(5.0), size=16)
    _notes(slide, items)


# ── public API ───────────────────────────────────────────────────────────────
def build_deck(meta: dict, analytics: dict, bundle: dict, narrative: dict) -> io.BytesIO:
    """Assemble the full reliability-review deck and return it as a BytesIO."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _cover(prs, meta)
    _slide_exec(prs, bundle, narrative)
    _slide_trends(prs, analytics)
    _slide_top_equipment(prs, analytics)
    _slide_rca(prs, analytics, narrative)
    _slide_capa(prs, bundle)
    _slide_risks(prs, bundle)
    _slide_actions(prs, narrative)
    _slide_decisions(prs, narrative)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
